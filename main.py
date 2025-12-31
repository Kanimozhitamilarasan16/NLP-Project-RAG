import os
import time
import threading
from PyPDF2 import PdfReader
import docx
from pptx import Presentation
from sentence_transformers import SentenceTransformer
import faiss
import tkinter as tk
from tkinter import scrolledtext, messagebox, ttk
from PIL import Image, ImageTk
from graphviz import Digraph
from transformers import pipeline

os.environ["PATH"] += os.pathsep + r"C:\Program Files\Graphviz\bin"

def load_docs(folder_path):
    texts = []
    for file in os.listdir(folder_path):
        path = os.path.join(folder_path, file)
        if file.lower().endswith(".pdf"):
            reader = PdfReader(path)
            text = ""
            for page in reader.pages:
                if page.extract_text():
                    text += page.extract_text() + "\n"
            texts.append(text)
        elif file.lower().endswith((".doc", ".docx")):
            doc = docx.Document(path)
            text = "\n".join([para.text for para in doc.paragraphs])
            texts.append(text)
        elif file.lower().endswith(".pptx"):
            prs = Presentation(path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            texts.append(text)
    return texts

def split_text(texts, chunk_size=500):
    chunks = []
    for text in texts:
        for i in range(0, len(text), chunk_size):
            chunks.append(text[i:i+chunk_size])
    return chunks

model = SentenceTransformer('all-MiniLM-L6-v2')
docs = load_docs("data/sample_docs")
chunks = split_text(docs)
embeddings = model.encode(chunks, convert_to_numpy=True)
dimension = embeddings.shape[1]
index = faiss.IndexFlatL2(dimension)
index.add(embeddings)

def query(text):
    query_emb = model.encode([text], convert_to_numpy=True)
    D, I = index.search(query_emb, k=3)
    context = "\n".join([chunks[i] for i in I[0]])
    return context

summarizer = pipeline("summarization", model="sshleifer/distilbart-cnn-12-6")

def generate_notes(text):
    notes = []
    for i in range(0, len(text), 500):
        chunk = text[i:i+500]
        summary = summarizer(chunk, max_length=100, min_length=30, do_sample=False)
        notes.append(summary[0]['summary_text'])
    return "\n".join(notes)

def generate_flowchart(text):
    timestamp = str(int(time.time() * 1000))
    output_file = f"flowchart_{timestamp}"
    sentences = [s.strip() for s in text.split('.') if s.strip()]
    dot = Digraph(comment='Flowchart', format='png')
    for i, sentence in enumerate(sentences):
        label = sentence if len(sentence) <= 50 else sentence[:50] + "..."
        dot.node(str(i), label)
        if i > 0:
            dot.edge(str(i-1), str(i))
    dot.render(output_file, cleanup=True)
    return output_file + ".png"

root = tk.Tk()
root.title("Offline RAG System - NLP Project")
root.geometry("1100x750")
root.configure(bg="#f0f4f7")

tk.Label(root, text="Intelligent Offline RAG System", font=("Helvetica", 22, "bold"), bg="#f0f4f7", fg="#2c3e50").pack(pady=15)

input_frame = tk.Frame(root, bg="#f0f4f7")
input_frame.pack(pady=10, fill=tk.X, padx=20)

tk.Label(input_frame, text="Enter your question:", font=("Helvetica", 14), bg="#f0f4f7").pack(side=tk.LEFT, padx=5)
question_entry = tk.Entry(input_frame, font=("Helvetica", 14), width=70)
question_entry.pack(side=tk.LEFT, padx=5)

progress_label = tk.Label(input_frame, text="", font=("Helvetica", 12), fg="#27ae60", bg="#f0f4f7")
progress_label.pack(side=tk.LEFT, padx=10)

tab_control = ttk.Notebook(root)
tab_control.pack(fill=tk.BOTH, expand=True, padx=20, pady=10)

answer_tab = ttk.Frame(tab_control)
tab_control.add(answer_tab, text="Answer")
output_text = scrolledtext.ScrolledText(answer_tab, font=("Helvetica", 12), wrap=tk.WORD, state='disabled', bg="#ffffff")
output_text.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)

notes_tab = ttk.Frame(tab_control)
tab_control.add(notes_tab, text="Notes & Flowchart")

tk.Label(notes_tab, text="Summarized Notes:", font=("Helvetica", 14)).pack(pady=5)
notes_text = scrolledtext.ScrolledText(notes_tab, font=("Helvetica", 12), wrap=tk.WORD, state='disabled', height=10, bg="#ffffff")
notes_text.pack(fill=tk.X, padx=10, pady=5)

flow_label = tk.Label(notes_tab)
flow_label.pack(pady=10)

def ask_question():
    question = question_entry.get().strip()
    if not question:
        messagebox.showwarning("Input Required", "Please enter a question!")
        return
    progress_label.config(text="Processing...")
    root.update_idletasks()
    try:
        result = query(question)
        output_text.configure(state='normal')
        output_text.delete(1.0, tk.END)
        output_text.insert(tk.END, result)
        output_text.configure(state='disabled')

        notes = generate_notes(result)
        notes_text.configure(state='normal')
        notes_text.delete(1.0, tk.END)
        notes_text.insert(tk.END, notes)
        notes_text.configure(state='disabled')

        img_path = generate_flowchart(result)
        img = Image.open(img_path)
        img = img.resize((800, 400), Image.Resampling.LANCZOS)
        img_tk = ImageTk.PhotoImage(img)
        flow_label.configure(image=img_tk)
        flow_label.image = img_tk
    except Exception as e:
        messagebox.showerror("Error", str(e))
    progress_label.config(text="Done!")

def ask_question_thread():
    threading.Thread(target=ask_question).start()

ask_btn = tk.Button(root, text="Ask Question", font=("Helvetica", 14, "bold"), bg="#2980b9", fg="white", command=ask_question_thread)
ask_btn.pack(pady=10)

root.mainloop()
